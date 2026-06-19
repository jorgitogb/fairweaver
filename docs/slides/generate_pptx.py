#!/usr/bin/env python3
"""Generate editable PowerPoint presentation for FAIRweaver workflow slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
DIAGRAM_DIR = os.path.join(os.path.dirname(__file__), "diagrams")

BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x61, 0x61, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_font(run, size=18, bold=False, color=BLACK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_title(slide, text, left=Inches(0.5), top=Inches(0.3), width=Inches(12), size=28, color=BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return txBox


def add_text_box(slide, text, left, top, width, height, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return txBox


def add_multiline_text(slide, lines, left, top, width, height, size=16, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_font(run, size=size, color=color)
    return txBox


def add_image(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        kwargs = {}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(img_path, left, top, **kwargs)


def add_table(slide, data, left, top, width, height, header_color=BLUE):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_font(run, size=14, bold=(i == 0), color=WHITE if i == 0 else BLACK)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
    return table_shape


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1 — Title
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Schema.org → ARC → FAIRagro Workflow",
              top=Inches(2.5), size=36, color=BLUE)
    add_text_box(slide, "TA3/4 Retreat 2026",
                 left=Inches(0.5), top=Inches(4), width=Inches(12), height=Inches(0.8), size=20, color=GRAY)
    add_notes(slide, (
        "Agricultural research data lives across dozens of repositories with different schemas, "
        "formats, and conventions. ARCs (Annotated Research Contexts) are the emerging standard for "
        "packaging that data, and FAIRagro's Search Hub is where researchers discover it. "
        "FAIRweaver's conversion pipeline bridges the gap: take metadata from any source, map it into "
        "an ARC, then extract the FAIRagro-specific schema for Search Hub ingestion.\n\n"
        "The challenge is structural. Every ARC is a JSON-LD graph, and different research groups "
        "model the same domain concepts at different depths and in different graph locations. "
        "The pipeline must handle that variability — or define conventions to eliminate it."
    ))

    # =========================================================================
    # SLIDE 2 — Pipeline
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "FAIRagro Metadata Transformation Pipeline")
    bullets = [
        "Sequential dependency: Output B (FAIRagro JSON) is derived from Output A (ARC RO-Crate) — not a parallel output",
        "Two harvest paths converge on the same FAIRagro JSON schema:",
        "  \u2022 Path 1 (solid): Direct harvest from GitLab DataHub where ARCs are stored",
        "  \u2022 Path 2 (dashed): Via FAIRagro Middleware API (federated service)",
        "Both paths produce identical FAIRagro schema.json ingested into SearchHub"
    ]
    add_multiline_text(slide, bullets, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    add_notes(slide, (
        "The pipeline is strictly sequential. Schema.org JSON-LD from an RDI (Research Data Infrastructure) "
        "enters at the top, gets transformed by FAIRagro templates, and produces two outputs in order:\n\n"
        "1. Output A — ARC RO-Crate. A JSON-LD document that follows the ARC specification.\n"
        "2. Output B — FAIRagro JSON. Derived from Output A, not produced independently.\n\n"
        "Output B's dependency on Output A is the key architectural insight. If the ARC is "
        "incomplete, the FAIRagro extraction is incomplete too. There is no shortcut.\n\n"
        "Two paths feed the ARC into Search Hub:\n"
        "- Path 1 (direct harvest): GitLab DataHub stores the ARC; a harvester reads it and produces FAIRagro JSON.\n"
        "- Path 2 (orchestrated): The FAIRagro Middleware API coordinates the full workflow, "
        "from harvest through conversion to ingestion.\n\n"
        "Both paths converge on the same FAIRagro schema.json consumed by Search Hub."
    ))

    # =========================================================================
    # SLIDE 3 — Three Scenarios
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Three File Scenarios: Input \u2192 ARC \u2192 FAIRagro Output")
    table_data = [
        ["Case", "Input File", "ARC Output", "FAIRagro Output"],
        ["Synthetic", "schema-org-wheat-full.json", "arc-ro-crate-wheat-full \u2713 compliant", "Full extraction \u2713"],
        ["Real \u2014 Small", "arc-ro-crate-dronflyover.json (<10 MB)", "Manual, partial \u26a0", "Partial \u2014 mappable fields only"],
        ["Real \u2014 Large", "arc-ro-crate-muenchenberg-lte.json (>100 MB)", "Manual, partial \u26a0", "Basic harvest only"]
    ]
    add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(12), Inches(2.5))
    add_text_box(slide,
                 "\U0001f4a1 If an ARC follows the FAIRagro specification \u2192 full metadata extraction. If not \u2192 only basic information is harvested.",
                 left=Inches(0.5), top=Inches(5), width=Inches(12), height=Inches(0.8), size=18, bold=True, color=GREEN)
    add_notes(slide, (
        "Three test files illustrate the relationship between ARC quality and extraction depth:\n\n"
        "Synthetic (schema-org-wheat-full.json): A manually crafted input designed to cover "
        "every FAIRagro-required field. The resulting ARC is fully compliant; extraction gets "
        "everything. This proves the pipeline works when the input follows the rules.\n\n"
        "Real \u2014 Small (arc-ro-crate-dronflyover.json, <10 MB): An actual agronomic ARC from "
        "the UC13 drone flyover. It follows the ISA Investigation \u2192 Study \u2192 Assay hierarchy but "
        "uses non-standard property names and missing optional fields. Result: partial extraction. "
        "Only fields that map unambiguously to the FAIRagro schema come through.\n\n"
        "Real \u2014 Large (arc-ro-crate-muenchenberg-lte.json, >100 MB): A long-term experiment "
        "ARC from M\u00fcncheberg with 27+ assays but a flatter structure. Many domain-specific fields "
        "don\u2019t match any FAIRagro template key. Result: basic harvest only \u2014 title, description, "
        "identifier, creator.\n\n"
        "The pattern is clear: compliance determines extraction depth. If an ARC follows the "
        "FAIRagro specification, you get full extraction. If not, you get the intersection of "
        "what\u2019s present and what\u2019s mappable."
    ))

    # =========================================================================
    # SLIDE 4 — Drone Structure
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Examining ARC Structure: Domain Objects at Different Depths")
    bullets = [
        "Goal:",
        "  \u2022 Understand how Agrischemas concepts map into ARC RO-Crate",
        "  \u2022 Show that equivalent domain concepts require very different traversal depths"
    ]
    add_multiline_text(slide, bullets, Inches(0.5), Inches(1.3), Inches(12), Inches(1.5))
    add_image(slide, os.path.join(DIAGRAM_DIR, "slide3-drone.png"),
              left=Inches(1.5), top=Inches(2.8), width=Inches(10))
    add_text_box(slide, "Example ARC RO-Crate: UC13 drone-flyover",
                 left=Inches(0.5), top=Inches(6.8), width=Inches(12), height=Inches(0.5), size=14, color=GRAY)
    add_notes(slide, (
        "Two real ARCs, one problem: same domain concepts at different graph depths.\n\n"
        "Drone Flyover (UC13):\n"
        "Following the ISA hierarchy: Investigation \u2192 Study \u2192 Assay. Two Agrischemas concepts live "
        "at very different traversal depths:\n\n"
        "- sensorType is 1 hop from Assay via measurementMethod \u2192 DefinedTerm.\n"
        "- cropSpecies is 4 hops from Study via about \u2192 LabProcess \u2192 object \u2192 Sample \u2192 "
        "additionalProperty \u2192 PropertyValue.\n\n"
        "The sensorType path is short and direct. The cropSpecies path requires traversing five "
        "nodes with specific predicates at each step. The parser must know the exact path \u2014 any "
        "deviation and it loses the concept."
    ))

    # =========================================================================
    # SLIDE 5 — Müncheberg
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "M\u00fcncheberg ARC: A Different Structural Pattern")
    bullets = [
        "Goal:",
        "  \u2022 Show another real ARC with a different structural pattern",
        "  \u2022 Reinforce that parser must handle multiple modeling conventions"
    ]
    add_multiline_text(slide, bullets, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2))
    add_image(slide, os.path.join(DIAGRAM_DIR, "slide4-muenchenberg.png"),
              left=Inches(3.5), top=Inches(2.3), width=Inches(9))
    add_text_box(slide, "Example ARC RO-Crate: M\u00fcncheberg LTE",
                 left=Inches(0.5), top=Inches(6.8), width=Inches(12), height=Inches(0.5), size=14, color=GRAY)

    # Comparison table on same slide
    table_data = [
        ["Aspect", "Drone Flyover", "M\u00fcncheberg LTE"],
        ["Study entity", "Explicit, in hasPart chain", "Present but disconnected (not in hasPart)"],
        ["Crop species path (short)", "Study \u2192 LabProcess \u2192 Sample \u2192 PropertyValue (4 hops)", "Source \u2192 additionalProperty \u2192 CharacteristicValue (2 hops)"],
        ["Crop species path (long)", "Same as short (only path)", "ALSO via LabProcess \u2192 object \u2192 Source \u2192 additionalProperty"],
        ["Sensor metadata", "Present (DefinedTerm)", "Absent"],
        ["Assay count", "1", "27+"]
    ]
    add_table(slide, table_data, Inches(0.5), Inches(3.5), Inches(12), Inches(3))
    add_notes(slide, (
        "M\u00fcncheberg LTE:\n"
        "A structurally different ARC. Investigation connects directly to 27+ Assays; a Study entity "
        "exists (studies/LTE-V140-Muencheberg/) but is NOT linked via the Investigation\u2019s hasPart "
        "chain. Crop species surfaces via a shorter alternative path: Source \u2192 "
        "additionalProperty \u2192 CharacteristicValue (2 hops instead of 4). M\u00fcncheberg also has "
        "the same long LabProcess chain as the drone flyover (via Study.about \u2192 LabProcess \u2192 "
        "object \u2192 Source \u2192 additionalProperty), so both paths coexist. Sensor metadata is "
        "absent because this ARC is about crop phenology, not remote sensing.\n\n"
        "A generic parser cannot reliably extract crop species from both ARCs without knowing "
        "which structural pattern to apply. This variability is the core motivation for standardization."
    ))

    # =========================================================================
    # SLIDE 6 — Required Pattern & Standardization Gap
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Required Modeling Pattern & Standardization Gap")
    bullets = [
        "Goal:",
        "  \u2022 Define the required path for unambiguous extraction",
        "  \u2022 Identify what still needs standardization"
    ]
    add_multiline_text(slide, bullets, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2))
    add_image(slide, os.path.join(DIAGRAM_DIR, "slide5-required.png"),
              left=Inches(3.5), top=Inches(2.3), width=Inches(8))
    add_text_box(slide, "In bold: required objects/properties to represent Crop",
                 left=Inches(0.5), top=Inches(5.8), width=Inches(12), height=Inches(0.5), size=14, color=GREEN)
    add_text_box(slide, "Example ARC RO-Crate: UC13 drone-flyover",
                 left=Inches(0.5), top=Inches(6.1), width=Inches(12), height=Inches(0.5), size=14, color=GRAY)

    # Open Questions table on same slide
    table_data = [
        ["Question", "Details"],
        ["Structure: ?", "How to formally specify the required traversal path?"],
        ["propertyID: SSSOM mapping", "How to standardize ontology term mappings?"]
    ]
    add_table(slide, table_data, Inches(0.5), Inches(6.4), Inches(12), Inches(0.9))

    add_notes(slide, (
        "After examining both ARCs, a required traversal path emerges for unambiguous crop extraction:\n\n"
        "Study \u2192 about \u2192 LabProcess \u2192 object \u2192 Sample/additionalType:Material\n"
        "  \u2192 additionalProperty \u2192 PropertyValue/propertyID\n\n"
        "The terminal node \u2014 PropertyValue with a propertyID \u2014 is critical. name: Organism and "
        "value: Solanum tuberosum are human-readable but machine-ambiguous. The propertyID "
        "linking to an ontology term (e.g., agrovoc:c_49904) makes it semantically precise and "
        "machine-resolvable.\n\n"
        "Two open questions remain:\n\n"
        "1. Structure specification. How do we formally specify a required graph traversal path "
        "so that tools can validate compliance automatically? A shape expression (ShEx) or "
        "SHACL shape? A simple property-path DSL in the pivot registry YAML?\n\n"
        "2. Ontology mapping standardization. propertyID values need a shared vocabulary of "
        "mappings so that every ARC uses the same ontology term for \u201ccrop species.\u201d SSSOM "
        "(Simple Standard for Sharing Ontological Mappings) is a candidate \u2014 but it needs "
        "community adoption.\n\n"
        "These are not FAIRweaver problems. They are community-level standardization problems that "
        "FAIRweaver exposes and quantifies."
    ))

    output_path = os.path.join(os.path.dirname(__file__), "fairst-workflow-slides-editable.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
